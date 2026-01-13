"""
Document Conversion Module
Handles conversions between PDF, Word, PowerPoint, Excel, and image formats.
"""

from __future__ import annotations

import io
import os
import tempfile
import zipfile
from typing import Optional, Literal
from pathlib import Path

from PIL import Image
from pypdf import PdfReader, PdfWriter
from reportlab.lib.pagesizes import letter, A4
from reportlab.pdfgen import canvas
from reportlab.lib.units import inch
import img2pdf

try:
    from pdf2image import convert_from_bytes
    PDF2IMAGE_AVAILABLE = True
except ImportError:
    PDF2IMAGE_AVAILABLE = False

try:
    from docx import Document
    from docx.shared import Inches, Pt
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches as PPTXInches, Pt as PPTXPt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    import pdfplumber
    PDFPLUMBER_AVAILABLE = True
except ImportError:
    PDFPLUMBER_AVAILABLE = False


class ConversionError(Exception):
    """Base exception for conversion operations."""
    pass


class PDFToImageConverter:
    """Convert PDF pages to images."""
    
    @staticmethod
    def convert(
        buffer: bytes,
        format: Literal["jpg", "png", "webp", "tiff", "bmp"] = "jpg",
        dpi: int = 150,
        pages: Optional[list[int]] = None
    ) -> list[bytes]:
        """
        Convert PDF pages to images.
        
        Args:
            buffer: PDF file bytes
            format: Output format (jpg, png, webp, tiff, bmp)
            dpi: Resolution in dots per inch
            pages: Specific page numbers (0-indexed), None for all pages
            
        Returns:
            List of image bytes for each page
        """
        # Try pdf2image first (better quality), fall back to PyMuPDF
        try:
            if PDF2IMAGE_AVAILABLE:
                return PDFToImageConverter._convert_with_pdf2image(buffer, format, dpi, pages)
        except Exception:
            pass
        
        # Fallback to PyMuPDF (works without poppler)
        return PDFToImageConverter._convert_with_pymupdf(buffer, format, dpi, pages)
    
    @staticmethod
    def _convert_with_pdf2image(
        buffer: bytes,
        format: Literal["jpg", "png", "webp", "tiff", "bmp"],
        dpi: int,
        pages: Optional[list[int]]
    ) -> list[bytes]:
        """Convert using pdf2image (requires poppler)."""
        try:
            if pages:
                first_page = min(pages) + 1
                last_page = max(pages) + 1
                images = convert_from_bytes(buffer, dpi=dpi, first_page=first_page, last_page=last_page)
            else:
                images = convert_from_bytes(buffer, dpi=dpi)
            
            results = []
            for i, img in enumerate(images):
                if pages and (i + (pages[0] if pages else 0)) not in pages:
                    continue
                    
                output = io.BytesIO()
                if format == "jpg":
                    img = img.convert("RGB")
                    img.save(output, format="JPEG", quality=85, optimize=True)
                elif format == "png":
                    img.save(output, format="PNG", optimize=True)
                elif format == "webp":
                    img.save(output, format="WEBP", quality=85, method=6)
                elif format == "tiff":
                    img.save(output, format="TIFF", compression="tiff_lzw")
                elif format == "bmp":
                    img.save(output, format="BMP")
                results.append(output.getvalue())
            
            return results
        except Exception as e:
            raise ConversionError(f"pdf2image conversion failed: {e}")
    
    @staticmethod
    def _convert_with_pymupdf(
        buffer: bytes,
        format: Literal["jpg", "png", "webp", "tiff", "bmp"],
        dpi: int,
        pages: Optional[list[int]]
    ) -> list[bytes]:
        """Convert using PyMuPDF/fitz (no poppler required)."""
        try:
            import fitz
            
            doc = fitz.open(stream=buffer, filetype="pdf")
            results = []
            
            zoom = dpi / 72.0
            mat = fitz.Matrix(zoom, zoom)
            
            page_range = pages if pages else range(len(doc))
            
            for page_num in page_range:
                if page_num >= len(doc):
                    continue
                    
                page = doc[page_num]
                pix = page.get_pixmap(matrix=mat, alpha=False)
                
                # PyMuPDF natively supports jpg and png
                if format == "jpg":
                    img_bytes = pix.tobytes("jpeg", jpg_quality=85)
                elif format == "png":
                    img_bytes = pix.tobytes("png")
                else:
                    # For other formats, use PIL
                    pil_img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    output = io.BytesIO()
                    
                    if format == "webp":
                        pil_img.save(output, format="WEBP", quality=85, method=6)
                    elif format == "tiff":
                        pil_img.save(output, format="TIFF", compression="tiff_lzw")
                    elif format == "bmp":
                        pil_img.save(output, format="BMP")
                    
                    img_bytes = output.getvalue()
                
                results.append(img_bytes)
            
            doc.close()
            return results
            
        except Exception as e:
            raise ConversionError(f"PDF to image conversion failed: {e}")
    
    @staticmethod
    def convert_to_zip(
        buffer: bytes,
        format: Literal["jpg", "png", "webp", "tiff", "bmp"] = "jpg",
        dpi: int = 150
    ) -> bytes:
        """Convert all PDF pages to images and return as ZIP file."""
        images = PDFToImageConverter.convert(buffer, format, dpi)
        
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zf:
            for i, img_bytes in enumerate(images):
                filename = f"page_{i + 1:03d}.{format}"
                zf.writestr(filename, img_bytes)
        
        return zip_buffer.getvalue()


class ImageToPDFConverter:
    """Convert images to PDF."""
    
    @staticmethod
    def convert(
        image_buffers: list[bytes],
        page_size: Literal["a4", "letter", "fit"] = "a4",
        margin: float = 0.5  # inches
    ) -> bytes:
        """
        Convert multiple images to a single PDF.
        
        Args:
            image_buffers: List of image file bytes
            page_size: Page size or 'fit' to match image dimensions
            margin: Page margin in inches
            
        Returns:
            PDF file bytes
        """
        if not image_buffers:
            raise ConversionError("No images provided")
        
        try:
            if page_size == "fit":
                # Use img2pdf for exact image dimensions
                return img2pdf.convert(image_buffers)
            
            # Create PDF with specified page size
            page_sizes = {
                "a4": A4,
                "letter": letter,
            }
            size = page_sizes.get(page_size, A4)
            
            output = io.BytesIO()
            c = canvas.Canvas(output, pagesize=size)
            page_width, page_height = size
            margin_pts = margin * inch
            
            for img_bytes in image_buffers:
                img = Image.open(io.BytesIO(img_bytes))
                
                # Calculate available space
                available_width = page_width - (2 * margin_pts)
                available_height = page_height - (2 * margin_pts)
                
                # Scale image to fit
                img_width, img_height = img.size
                scale = min(available_width / img_width, available_height / img_height)
                
                new_width = img_width * scale
                new_height = img_height * scale
                
                # Center image
                x = margin_pts + (available_width - new_width) / 2
                y = margin_pts + (available_height - new_height) / 2
                
                # Save temp image for reportlab
                temp_img = io.BytesIO()
                if img.mode in ('RGBA', 'P'):
                    img = img.convert('RGB')
                img.save(temp_img, format='JPEG', quality=95)
                temp_img.seek(0)
                
                from reportlab.lib.utils import ImageReader
                c.drawImage(ImageReader(temp_img), x, y, width=new_width, height=new_height)
                c.showPage()
            
            c.save()
            return output.getvalue()
            
        except Exception as e:
            raise ConversionError(f"Image to PDF conversion failed: {e}")


class PDFToWordConverter:
    """Convert PDF to Word document."""
    
    @staticmethod
    def convert(buffer: bytes) -> bytes:
        """
        Convert PDF to DOCX format.
        Note: This creates a basic text extraction. Complex layouts may not be preserved.
        """
        if not DOCX_AVAILABLE:
            raise ConversionError("python-docx library not available")
        
        if not PDFPLUMBER_AVAILABLE:
            raise ConversionError("pdfplumber library not available")
        
        try:
            doc = Document()
            
            with pdfplumber.open(io.BytesIO(buffer)) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    # Extract text
                    text = page.extract_text()
                    if text:
                        for paragraph in text.split('\n\n'):
                            if paragraph.strip():
                                doc.add_paragraph(paragraph.strip())
                    
                    # Add page break between pages (except last)
                    if page_num < len(pdf.pages) - 1:
                        doc.add_page_break()
            
            output = io.BytesIO()
            doc.save(output)
            return output.getvalue()
            
        except Exception as e:
            raise ConversionError(f"PDF to Word conversion failed: {e}")


class WordToPDFConverter:
    """Convert Word documents to PDF."""
    
    @staticmethod
    def convert(buffer: bytes) -> bytes:
        """
        Convert DOCX to PDF format.
        Uses reportlab for cross-platform compatibility.
        """
        if not DOCX_AVAILABLE:
            raise ConversionError("python-docx library not available")
        
        try:
            doc = Document(io.BytesIO(buffer))
            
            output = io.BytesIO()
            c = canvas.Canvas(output, pagesize=letter)
            page_width, page_height = letter
            
            margin = 72  # 1 inch
            y_position = page_height - margin
            line_height = 14
            
            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    y_position -= line_height
                    continue
                
                # Simple word wrapping
                words = text.split()
                lines = []
                current_line = []
                
                for word in words:
                    current_line.append(word)
                    test_line = ' '.join(current_line)
                    text_width = c.stringWidth(test_line, "Helvetica", 11)
                    
                    if text_width > (page_width - 2 * margin):
                        current_line.pop()
                        if current_line:
                            lines.append(' '.join(current_line))
                        current_line = [word]
                
                if current_line:
                    lines.append(' '.join(current_line))
                
                # Check for page break
                if y_position - (len(lines) * line_height) < margin:
                    c.showPage()
                    y_position = page_height - margin
                
                for line in lines:
                    c.setFont("Helvetica", 11)
                    c.drawString(margin, y_position, line)
                    y_position -= line_height
                
                y_position -= line_height / 2  # Paragraph spacing
            
            c.save()
            return output.getvalue()
            
        except Exception as e:
            raise ConversionError(f"Word to PDF conversion failed: {e}")


class PDFToExcelConverter:
    """Convert PDF tables to Excel."""
    
    @staticmethod
    def convert(buffer: bytes) -> bytes:
        """
        Extract tables from PDF and convert to Excel.
        """
        if not OPENPYXL_AVAILABLE:
            raise ConversionError("openpyxl library not available")
        
        if not PDFPLUMBER_AVAILABLE:
            raise ConversionError("pdfplumber library not available")
        
        try:
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = "Extracted Data"
            
            row_offset = 1
            
            with pdfplumber.open(io.BytesIO(buffer)) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    tables = page.extract_tables()
                    
                    if tables:
                        for table in tables:
                            for row in table:
                                for col_idx, cell in enumerate(row):
                                    ws.cell(
                                        row=row_offset, 
                                        column=col_idx + 1, 
                                        value=cell if cell else ""
                                    )
                                row_offset += 1
                            row_offset += 1  # Space between tables
                    else:
                        # If no tables found, extract text
                        text = page.extract_text()
                        if text:
                            ws.cell(row=row_offset, column=1, value=f"--- Page {page_num + 1} ---")
                            row_offset += 1
                            for line in text.split('\n'):
                                ws.cell(row=row_offset, column=1, value=line)
                                row_offset += 1
                            row_offset += 1
            
            output = io.BytesIO()
            wb.save(output)
            return output.getvalue()
            
        except Exception as e:
            raise ConversionError(f"PDF to Excel conversion failed: {e}")


class ExcelToPDFConverter:
    """Convert Excel to PDF."""
    
    @staticmethod
    def convert(buffer: bytes) -> bytes:
        """Convert Excel file to PDF format with better formatting."""
        if not OPENPYXL_AVAILABLE:
            raise ConversionError("openpyxl library not available")
        
        try:
            wb = openpyxl.load_workbook(io.BytesIO(buffer), data_only=True)
            
            output = io.BytesIO()
            c = canvas.Canvas(output, pagesize=letter)
            page_width, page_height = letter
            
            margin = 40
            
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                
                # Calculate column widths based on content
                col_widths = {}
                max_row = 0
                for row_idx, row in enumerate(ws.iter_rows(values_only=True), 1):
                    max_row = max(max_row, row_idx)
                    for col_idx, cell_value in enumerate(row):
                        if cell_value is not None:
                            width = len(str(cell_value)) * 6 + 10
                            col_widths[col_idx] = max(col_widths.get(col_idx, 60), min(width, 120))
                
                # Calculate total width needed
                total_width = sum(col_widths.values())
                available_width = page_width - (2 * margin)
                
                # Scale if needed
                if total_width > available_width:
                    scale_factor = available_width / total_width
                    col_widths = {k: v * scale_factor for k, v in col_widths.items()}
                
                y_position = page_height - margin
                row_height = 16
                
                # Sheet title
                c.setFont("Helvetica-Bold", 12)
                c.drawString(margin, y_position, f"Sheet: {sheet_name}")
                y_position -= row_height * 1.5
                
                # Draw table
                for row_idx, row in enumerate(ws.iter_rows(values_only=True)):
                    if y_position < margin + row_height * 2:
                        c.showPage()
                        y_position = page_height - margin
                    
                    # Use bold for first row (headers)
                    if row_idx == 0:
                        c.setFont("Helvetica-Bold", 9)
                    else:
                        c.setFont("Helvetica", 9)
                    
                    x_position = margin
                    for col_idx, cell_value in enumerate(row):
                        if x_position > page_width - margin:
                            break
                            
                        if cell_value is not None:
                            # Handle numeric formatting
                            if isinstance(cell_value, (int, float)):
                                if isinstance(cell_value, float):
                                    text = f"{cell_value:.2f}"
                                else:
                                    text = str(cell_value)
                            else:
                                text = str(cell_value)
                            
                            # Truncate if too long
                            col_w = col_widths.get(col_idx, 60)
                            max_chars = int(col_w / 6)
                            if len(text) > max_chars:
                                text = text[:max_chars-2] + ".."
                            
                            c.drawString(x_position, y_position, text)
                        
                        x_position += col_widths.get(col_idx, 60)
                    
                    y_position -= row_height
                
                # New page for each sheet
                c.showPage()
                y_position = page_height - margin
            
            c.save()
            return output.getvalue()
            
        except Exception as e:
            raise ConversionError(f"Excel to PDF conversion failed: {e}")


class PDFToPowerPointConverter:
    """Convert PDF to PowerPoint."""
    
    @staticmethod
    def convert(buffer: bytes, dpi: int = 150) -> bytes:
        """
        Convert PDF to PowerPoint by converting each page to an image slide.
        """
        if not PPTX_AVAILABLE:
            raise ConversionError("python-pptx library not available")
        
        try:
            import fitz
            
            # Use PyMuPDF to convert pages to images
            doc = fitz.open(stream=buffer, filetype="pdf")
            
            # Create PowerPoint presentation
            prs = Presentation()
            prs.slide_width = PPTXInches(13.333)  # 16:9 widescreen
            prs.slide_height = PPTXInches(7.5)
            
            blank_layout = prs.slide_layouts[6]  # Blank layout
            
            zoom = dpi / 72.0
            mat = fitz.Matrix(zoom, zoom)
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                pix = page.get_pixmap(matrix=mat, alpha=False)
                
                # Convert pixmap to PNG bytes
                img_bytes = pix.tobytes("png")
                img_buffer = io.BytesIO(img_bytes)
                
                # Add slide
                slide = prs.slides.add_slide(blank_layout)
                
                # Add image to fill slide
                slide.shapes.add_picture(
                    img_buffer,
                    PPTXInches(0.25),
                    PPTXInches(0.25),
                    width=prs.slide_width - PPTXInches(0.5),
                    height=prs.slide_height - PPTXInches(0.5)
                )
            
            doc.close()
            
            output = io.BytesIO()
            prs.save(output)
            return output.getvalue()
            
        except Exception as e:
            raise ConversionError(f"PDF to PowerPoint conversion failed: {e}")


class PowerPointToPDFConverter:
    """Convert PowerPoint to PDF."""
    
    @staticmethod
    def convert(buffer: bytes) -> bytes:
        """
        Convert PowerPoint to PDF using the best available method.
        
        Strategy:
        1. Try LibreOffice first (free, unlimited)
        2. If any blank pages detected, use HYBRID approach:
           - Keep good pages from LibreOffice
           - Re-render blank pages using slide-to-image method
        3. CloudConvert as final fallback (limited quota)
        
        Set PPTX_CONVERT_METHOD env var to force: 'libreoffice', 'cloudconvert', 'hybrid', 'sliderender'
        """
        import platform
        import os
        
        forced_method = os.environ.get('PPTX_CONVERT_METHOD', '').lower()
        print(f"[PPTX→PDF] Starting conversion, forced_method={forced_method or 'auto'}")
        
        # Forced method handling
        if forced_method == 'cloudconvert':
            cloudconvert_key = os.environ.get('CLOUDCONVERT_API_KEY')
            if cloudconvert_key:
                return PowerPointToPDFConverter._convert_with_cloudconvert(buffer, cloudconvert_key)
            raise ConversionError("CLOUDCONVERT_API_KEY not set")
        elif forced_method == 'sliderender':
            return PowerPointToPDFConverter._convert_slides_to_images_pdf(buffer)
        elif forced_method == 'libreoffice':
            return PowerPointToPDFConverter._convert_with_libreoffice(buffer, strict=False)
        elif forced_method == 'hybrid':
            return PowerPointToPDFConverter._convert_hybrid(buffer)
        
        # Auto mode: Try LibreOffice, then hybrid fix, then CloudConvert
        cloudconvert_key = os.environ.get('CLOUDCONVERT_API_KEY')
        
        try:
            print("[PPTX→PDF] Trying LibreOffice...")
            result, blank_pages = PowerPointToPDFConverter._convert_with_libreoffice(buffer, strict=False, return_blank_info=True)
            
            if not blank_pages:
                print("[PPTX→PDF] LibreOffice succeeded with all pages valid!")
                return result
            
            # LibreOffice produced some blank pages - try hybrid fix
            print(f"[PPTX→PDF] LibreOffice produced {len(blank_pages)} blank pages: {blank_pages}")
            print("[PPTX→PDF] Attempting hybrid fix (re-render blank pages)...")
            
            try:
                fixed_result = PowerPointToPDFConverter._fix_blank_pages(buffer, result, blank_pages)
                print("[PPTX→PDF] Hybrid fix succeeded!")
                return fixed_result
            except Exception as fix_error:
                print(f"[PPTX→PDF] Hybrid fix failed: {fix_error}")
                
        except Exception as e:
            print(f"[PPTX→PDF] LibreOffice error: {e}")
        
        # Try pure slide-to-image rendering
        try:
            print("[PPTX→PDF] Trying slide-to-image rendering...")
            result = PowerPointToPDFConverter._convert_slides_to_images_pdf(buffer)
            print("[PPTX→PDF] Slide rendering succeeded!")
            return result
        except Exception as e:
            print(f"[PPTX→PDF] Slide rendering failed: {e}")
        
        # Windows: try PowerPoint COM
        if platform.system() == "Windows":
            try:
                print("[PPTX→PDF] Trying PowerPoint COM...")
                result = PowerPointToPDFConverter._convert_with_com(buffer)
                print("[PPTX→PDF] COM succeeded!")
                return result
            except Exception as e:
                print(f"[PPTX→PDF] COM failed: {e}")
        
        # CloudConvert as final fallback
        if cloudconvert_key:
            try:
                print("[PPTX→PDF] Final attempt with CloudConvert...")
                return PowerPointToPDFConverter._convert_with_cloudconvert(buffer, cloudconvert_key)
            except Exception as e:
                print(f"[PPTX→PDF] CloudConvert failed: {e}")
        
        raise ConversionError(
            "Could not convert PPTX to PDF. All methods failed."
        )
    
    @staticmethod
    def _convert_slides_to_images_pdf(buffer: bytes) -> bytes:
        """
        Convert PPTX by rendering each slide to a high-quality image, then combining into PDF.
        This method extracts ALL visual content from slides.
        """
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from PIL import Image, ImageDraw, ImageFont
        import fitz  # PyMuPDF
        
        print("[SlideRender] Starting slide-to-image conversion...")
        
        prs = Presentation(io.BytesIO(buffer))
        slide_width_pt = prs.slide_width.pt
        slide_height_pt = prs.slide_height.pt
        
        # Render at 2x resolution for quality (150 DPI equivalent)
        scale = 2.0
        img_width = int(slide_width_pt * scale)
        img_height = int(slide_height_pt * scale)
        
        print(f"[SlideRender] Slide size: {slide_width_pt}x{slide_height_pt}pt, rendering at {img_width}x{img_height}px")
        
        slide_images = []
        
        for slide_idx, slide in enumerate(prs.slides):
            print(f"[SlideRender] Rendering slide {slide_idx + 1}/{len(prs.slides)}...")
            
            # Create blank white image
            img = Image.new('RGB', (img_width, img_height), (255, 255, 255))
            draw = ImageDraw.Draw(img)
            
            # Try to render slide background
            try:
                bg = slide.background
                if bg and bg.fill and bg.fill.type is not None:
                    if hasattr(bg.fill, 'fore_color') and bg.fill.fore_color:
                        try:
                            rgb = bg.fill.fore_color.rgb
                            if rgb:
                                draw.rectangle([0, 0, img_width, img_height], 
                                             fill=(rgb[0], rgb[1], rgb[2]))
                        except:
                            pass
            except:
                pass
            
            # Collect shapes sorted by z-order (approximate using position)
            shapes_to_render = []
            for shape in slide.shapes:
                try:
                    x = shape.left.pt * scale if hasattr(shape, 'left') else 0
                    y = shape.top.pt * scale if hasattr(shape, 'top') else 0
                    w = shape.width.pt * scale if hasattr(shape, 'width') else 100
                    h = shape.height.pt * scale if hasattr(shape, 'height') else 100
                    shapes_to_render.append({
                        'shape': shape,
                        'x': x, 'y': y, 'w': w, 'h': h
                    })
                except:
                    continue
            
            # Render shapes (images first, then shapes, then text)
            for shape_info in shapes_to_render:
                shape = shape_info['shape']
                x, y, w, h = int(shape_info['x']), int(shape_info['y']), int(shape_info['w']), int(shape_info['h'])
                
                try:
                    # Render pictures/images
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            img_data = shape.image.blob
                            shape_img = Image.open(io.BytesIO(img_data))
                            
                            # Handle transparency
                            if shape_img.mode in ('RGBA', 'LA', 'P'):
                                if shape_img.mode == 'P':
                                    shape_img = shape_img.convert('RGBA')
                                # Composite onto white background
                                bg_img = Image.new('RGB', shape_img.size, (255, 255, 255))
                                if shape_img.mode == 'RGBA':
                                    bg_img.paste(shape_img, mask=shape_img.split()[3])
                                else:
                                    bg_img.paste(shape_img)
                                shape_img = bg_img
                            elif shape_img.mode != 'RGB':
                                shape_img = shape_img.convert('RGB')
                            
                            # Resize and paste
                            shape_img = shape_img.resize((max(1, w), max(1, h)), Image.Resampling.LANCZOS)
                            img.paste(shape_img, (x, y))
                        except Exception as e:
                            # Draw placeholder for failed images
                            draw.rectangle([x, y, x + w, y + h], outline=(200, 200, 200), fill=(245, 245, 245))
                        continue
                    
                    # Render shape fills
                    if hasattr(shape, 'fill') and shape.fill:
                        try:
                            fill = shape.fill
                            if fill.type is not None and hasattr(fill, 'fore_color') and fill.fore_color:
                                rgb = fill.fore_color.rgb
                                if rgb:
                                    draw.rectangle([x, y, x + w, y + h], fill=(rgb[0], rgb[1], rgb[2]))
                        except:
                            pass
                    
                    # Render shape borders
                    if hasattr(shape, 'line') and shape.line:
                        try:
                            line = shape.line
                            if line.color and line.color.rgb:
                                rgb = line.color.rgb
                                line_w = int((line.width.pt if line.width else 1) * scale)
                                draw.rectangle([x, y, x + w, y + h], outline=(rgb[0], rgb[1], rgb[2]), width=max(1, line_w))
                        except:
                            pass
                    
                    # Render text
                    if hasattr(shape, 'text_frame'):
                        try:
                            text_frame = shape.text_frame
                            text_y = y + 5
                            
                            for para in text_frame.paragraphs:
                                para_text = ""
                                text_color = (0, 0, 0)
                                font_size = 14
                                
                                for run in para.runs:
                                    para_text += run.text
                                    if run.font.size:
                                        font_size = int(run.font.size.pt * scale * 0.8)
                                    if run.font.color and run.font.color.rgb:
                                        rgb = run.font.color.rgb
                                        text_color = (rgb[0], rgb[1], rgb[2])
                                
                                if para_text.strip():
                                    # Try to use a font, fall back to default
                                    try:
                                        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", max(8, font_size))
                                    except:
                                        try:
                                            font = ImageFont.truetype("arial.ttf", max(8, font_size))
                                        except:
                                            font = ImageFont.load_default()
                                    
                                    # Word wrap
                                    max_width = w - 10
                                    words = para_text.split()
                                    lines = []
                                    current_line = []
                                    
                                    for word in words:
                                        current_line.append(word)
                                        test_line = ' '.join(current_line)
                                        bbox = draw.textbbox((0, 0), test_line, font=font)
                                        if bbox[2] - bbox[0] > max_width and len(current_line) > 1:
                                            current_line.pop()
                                            lines.append(' '.join(current_line))
                                            current_line = [word]
                                    if current_line:
                                        lines.append(' '.join(current_line))
                                    
                                    for line in lines:
                                        if text_y < y + h - 5:
                                            draw.text((x + 5, text_y), line, fill=text_color, font=font)
                                            text_y += font_size + 4
                        except:
                            pass
                            
                except Exception as shape_error:
                    continue
            
            slide_images.append(img)
        
        # Combine images into PDF using PyMuPDF
        print(f"[SlideRender] Combining {len(slide_images)} slides into PDF...")
        
        pdf_doc = fitz.open()
        
        for slide_img in slide_images:
            # Convert PIL image to bytes
            img_buffer = io.BytesIO()
            slide_img.save(img_buffer, format='PNG', optimize=True)
            img_buffer.seek(0)
            
            # Create PDF page with slide dimensions
            page = pdf_doc.new_page(width=slide_width_pt, height=slide_height_pt)
            
            # Insert image to fill page
            rect = fitz.Rect(0, 0, slide_width_pt, slide_height_pt)
            page.insert_image(rect, stream=img_buffer.getvalue())
        
        # Save PDF
        pdf_buffer = io.BytesIO()
        pdf_doc.save(pdf_buffer)
        pdf_doc.close()
        
        result = pdf_buffer.getvalue()
        print(f"[SlideRender] PDF created: {len(result)} bytes")
        return result
    
    @staticmethod
    def _fix_blank_pages(pptx_buffer: bytes, pdf_data: bytes, blank_pages: list) -> bytes:
        """
        Fix blank pages in LibreOffice output by re-rendering those specific pages
        from the original PPTX and merging them back.
        """
        import fitz
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from PIL import Image, ImageDraw, ImageFont
        
        print(f"[HybridFix] Fixing {len(blank_pages)} blank pages: {blank_pages}")
        
        # Open the PDF
        pdf_doc = fitz.open(stream=pdf_data, filetype="pdf")
        
        # Load PPTX
        prs = Presentation(io.BytesIO(pptx_buffer))
        slide_width_pt = prs.slide_width.pt
        slide_height_pt = prs.slide_height.pt
        
        scale = 2.0
        img_width = int(slide_width_pt * scale)
        img_height = int(slide_height_pt * scale)
        
        # Re-render each blank page
        for page_num in blank_pages:
            slide_idx = page_num - 1  # Convert to 0-indexed
            
            if slide_idx >= len(prs.slides):
                continue
            
            slide = prs.slides[slide_idx]
            print(f"[HybridFix] Re-rendering page {page_num}...")
            
            # Render slide to image (same logic as _convert_slides_to_images_pdf)
            img = Image.new('RGB', (img_width, img_height), (255, 255, 255))
            draw = ImageDraw.Draw(img)
            
            # Render background
            try:
                bg = slide.background
                if bg and bg.fill and bg.fill.type is not None:
                    if hasattr(bg.fill, 'fore_color') and bg.fill.fore_color:
                        try:
                            rgb = bg.fill.fore_color.rgb
                            if rgb:
                                draw.rectangle([0, 0, img_width, img_height], fill=(rgb[0], rgb[1], rgb[2]))
                        except:
                            pass
            except:
                pass
            
            # Render all shapes
            for shape in slide.shapes:
                try:
                    x = int(shape.left.pt * scale) if hasattr(shape, 'left') else 0
                    y = int(shape.top.pt * scale) if hasattr(shape, 'top') else 0
                    w = int(shape.width.pt * scale) if hasattr(shape, 'width') else 100
                    h = int(shape.height.pt * scale) if hasattr(shape, 'height') else 100
                    
                    # Pictures
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            shape_img = Image.open(io.BytesIO(shape.image.blob))
                            if shape_img.mode in ('RGBA', 'LA', 'P'):
                                if shape_img.mode == 'P':
                                    shape_img = shape_img.convert('RGBA')
                                bg_img = Image.new('RGB', shape_img.size, (255, 255, 255))
                                if shape_img.mode == 'RGBA':
                                    bg_img.paste(shape_img, mask=shape_img.split()[3])
                                else:
                                    bg_img.paste(shape_img)
                                shape_img = bg_img
                            elif shape_img.mode != 'RGB':
                                shape_img = shape_img.convert('RGB')
                            shape_img = shape_img.resize((max(1, w), max(1, h)), Image.Resampling.LANCZOS)
                            img.paste(shape_img, (x, y))
                        except:
                            draw.rectangle([x, y, x + w, y + h], outline=(200, 200, 200), fill=(245, 245, 245))
                        continue
                    
                    # Shape fills
                    if hasattr(shape, 'fill') and shape.fill:
                        try:
                            if shape.fill.type is not None and hasattr(shape.fill, 'fore_color') and shape.fill.fore_color:
                                rgb = shape.fill.fore_color.rgb
                                if rgb:
                                    draw.rectangle([x, y, x + w, y + h], fill=(rgb[0], rgb[1], rgb[2]))
                        except:
                            pass
                    
                    # Text
                    if hasattr(shape, 'text_frame'):
                        try:
                            text_y = y + 5
                            for para in shape.text_frame.paragraphs:
                                para_text = ""
                                text_color = (0, 0, 0)
                                font_size = int(14 * scale)
                                
                                for run in para.runs:
                                    para_text += run.text
                                    if run.font.size:
                                        font_size = int(run.font.size.pt * scale * 0.8)
                                    if run.font.color and run.font.color.rgb:
                                        rgb = run.font.color.rgb
                                        text_color = (rgb[0], rgb[1], rgb[2])
                                
                                if para_text.strip():
                                    try:
                                        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", max(8, font_size))
                                    except:
                                        font = ImageFont.load_default()
                                    
                                    draw.text((x + 5, text_y), para_text, fill=text_color, font=font)
                                    text_y += font_size + 4
                        except:
                            pass
                except:
                    continue
            
            # Convert image to PNG bytes
            img_buffer = io.BytesIO()
            img.save(img_buffer, format='PNG')
            img_buffer.seek(0)
            
            # Replace the blank page in PDF
            page = pdf_doc[slide_idx]
            
            # Clear the page and insert new image
            page.clean_contents()
            rect = fitz.Rect(0, 0, slide_width_pt, slide_height_pt)
            page.insert_image(rect, stream=img_buffer.getvalue())
        
        # Save modified PDF
        pdf_buffer = io.BytesIO()
        pdf_doc.save(pdf_buffer)
        pdf_doc.close()
        
        result = pdf_buffer.getvalue()
        print(f"[HybridFix] Fixed PDF: {len(result)} bytes")
        return result
    
    @staticmethod
    def _convert_hybrid(buffer: bytes) -> bytes:
        """Force hybrid mode: LibreOffice + fix any blank pages."""
        result, blank_pages = PowerPointToPDFConverter._convert_with_libreoffice(buffer, strict=False, return_blank_info=True)
        
        if blank_pages:
            return PowerPointToPDFConverter._fix_blank_pages(buffer, result, blank_pages)
        return result
    
    @staticmethod
    def _convert_with_cloudconvert(buffer: bytes, api_key: str) -> bytes:
        """Convert using CloudConvert API (25 free/day)."""
        import requests
        import time
        
        headers = {"Authorization": f"Bearer {api_key}"}
        base_url = "https://api.cloudconvert.com/v2"
        
        # Step 1: Create job
        job_response = requests.post(
            f"{base_url}/jobs",
            headers=headers,
            json={
                "tasks": {
                    "upload-file": {"operation": "import/upload"},
                    "convert-file": {
                        "operation": "convert",
                        "input": ["upload-file"],
                        "output_format": "pdf"
                    },
                    "export-file": {
                        "operation": "export/url",
                        "input": ["convert-file"]
                    }
                }
            }
        )
        job_response.raise_for_status()
        job_data = job_response.json()
        
        # Step 2: Upload file
        upload_task = next(t for t in job_data["data"]["tasks"] if t["name"] == "upload-file")
        upload_url = upload_task["result"]["form"]["url"]
        upload_params = upload_task["result"]["form"]["parameters"]
        
        files = {"file": ("input.pptx", buffer, "application/vnd.openxmlformats-officedocument.presentationml.presentation")}
        upload_response = requests.post(upload_url, data=upload_params, files=files)
        upload_response.raise_for_status()
        
        # Step 3: Wait for conversion
        job_id = job_data["data"]["id"]
        for _ in range(60):  # Wait up to 60 seconds
            time.sleep(1)
            status_response = requests.get(f"{base_url}/jobs/{job_id}", headers=headers)
            status_data = status_response.json()
            
            if status_data["data"]["status"] == "finished":
                # Get download URL
                export_task = next(t for t in status_data["data"]["tasks"] if t["name"] == "export-file")
                download_url = export_task["result"]["files"][0]["url"]
                
                # Download PDF
                pdf_response = requests.get(download_url)
                return pdf_response.content
            elif status_data["data"]["status"] == "error":
                raise ConversionError("CloudConvert conversion failed")
        
        raise ConversionError("CloudConvert conversion timed out")
    
    @staticmethod
    def _convert_with_com(buffer: bytes) -> bytes:
        """Convert using Microsoft PowerPoint COM (Windows only)."""
        import tempfile
        import os
        import time
        
        try:
            import comtypes.client
        except ImportError:
            raise ConversionError("comtypes not available")
        
        # Save to temp file with unique name
        temp_dir = tempfile.mkdtemp()
        pptx_path = os.path.join(temp_dir, 'input.pptx')
        pdf_path = os.path.join(temp_dir, 'output.pdf')
        
        powerpoint = None
        presentation = None
        
        try:
            with open(pptx_path, 'wb') as f:
                f.write(buffer)
            
            # Initialize COM
            comtypes.client.CoInitialize()
            
            powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
            powerpoint.Visible = 1  # Must be visible for proper rendering
            powerpoint.DisplayAlerts = 0  # Suppress dialogs
            
            # Open presentation with specific settings
            # ReadOnly=False, Untitled=False, WithWindow=True for proper rendering
            presentation = powerpoint.Presentations.Open(
                pptx_path,
                ReadOnly=True,
                Untitled=False,
                WithWindow=True
            )
            
            # Wait for presentation to fully load
            time.sleep(0.5)
            
            # Export to PDF with high quality settings
            # 32 = ppSaveAsPDF
            # Using ExportAsFixedFormat for more control
            try:
                # ppFixedFormatTypePDF = 2
                # ppFixedFormatIntentPrint = 2 (high quality)
                presentation.ExportAsFixedFormat(
                    pdf_path,
                    2,  # ppFixedFormatTypePDF
                    Intent=2,  # ppFixedFormatIntentPrint (high quality)
                    FrameSlides=False,
                    HandoutOrder=1,
                    OutputType=1,  # All slides
                    PrintHiddenSlides=False,
                    IncludeDocProperties=True,
                    DocStructureTags=True
                )
            except Exception:
                # Fallback to SaveAs if ExportAsFixedFormat fails
                presentation.SaveAs(pdf_path, 32)  # 32 = ppSaveAsPDF
            
            # Close properly
            presentation.Close()
            presentation = None
            
            powerpoint.Quit()
            powerpoint = None
            
            # Small delay to ensure file is written
            time.sleep(0.3)
            
            if not os.path.exists(pdf_path):
                raise ConversionError("PDF file was not created")
            
            with open(pdf_path, 'rb') as f:
                result = f.read()
            
            if len(result) < 1000:
                raise ConversionError("PDF file appears to be empty or corrupted")
            
            return result
            
        except Exception as e:
            raise ConversionError(f"COM conversion error: {e}")
        finally:
            # Cleanup COM objects
            try:
                if presentation:
                    presentation.Close()
            except:
                pass
            try:
                if powerpoint:
                    powerpoint.Quit()
            except:
                pass
            
            # Cleanup temp files
            time.sleep(0.2)
            try:
                if os.path.exists(pptx_path):
                    os.unlink(pptx_path)
                if os.path.exists(pdf_path):
                    os.unlink(pdf_path)
                os.rmdir(temp_dir)
            except:
                pass
    
    @staticmethod
    def _convert_with_libreoffice(buffer: bytes, strict: bool = True, return_blank_info: bool = False):
        """
        Convert using LibreOffice with optimized settings for accuracy.
        
        Args:
            buffer: PPTX file bytes
            strict: If True, fail if ANY blank page is detected (default True)
            return_blank_info: If True, return tuple (pdf_data, blank_pages_list)
        
        Returns:
            bytes: PDF data (if return_blank_info=False)
            tuple: (pdf_data, blank_pages_list) if return_blank_info=True
        """
        import tempfile
        import subprocess
        import os
        import fitz  # PyMuPDF for PDF validation
        
        with tempfile.TemporaryDirectory() as tmpdir:
            pptx_path = os.path.join(tmpdir, 'input.pptx')
            
            with open(pptx_path, 'wb') as f:
                f.write(buffer)
            
            # Count slides and analyze content in PPTX for validation
            expected_pages = 0
            slide_contents = []  # Track what each slide contains
            try:
                from pptx import Presentation
                from pptx.enum.shapes import MSO_SHAPE_TYPE
                prs = Presentation(io.BytesIO(buffer))
                expected_pages = len(prs.slides)
                
                for slide_idx, slide in enumerate(prs.slides):
                    has_text = False
                    has_image = False
                    has_shape = False
                    
                    for shape in slide.shapes:
                        if hasattr(shape, 'text') and shape.text.strip():
                            has_text = True
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            has_image = True
                        if hasattr(shape, 'fill') or hasattr(shape, 'line'):
                            has_shape = True
                    
                    slide_contents.append({
                        'has_text': has_text,
                        'has_image': has_image,
                        'has_shape': has_shape,
                        'has_content': has_text or has_image or has_shape
                    })
                
                print(f"[LibreOffice] Input PPTX: {expected_pages} slides")
                content_slides = sum(1 for s in slide_contents if s['has_content'])
                print(f"[LibreOffice] Slides with content: {content_slides}")
            except Exception as e:
                print(f"[LibreOffice] Could not analyze PPTX: {e}")
                expected_pages = 0
            
            # LibreOffice paths to try
            soffice_paths = [
                '/usr/bin/soffice',
                '/usr/bin/libreoffice',
                'soffice',
                'libreoffice',
                'C:\\Program Files\\LibreOffice\\program\\soffice.exe',
                'C:\\Program Files (x86)\\LibreOffice\\program\\soffice.exe',
            ]
            
            # Set environment for better LibreOffice rendering
            env = os.environ.copy()
            env['HOME'] = tmpdir  # Avoid profile lock issues
            env['SAL_USE_VCLPLUGIN'] = 'svp'  # Use headless rendering
            
            for soffice in soffice_paths:
                try:
                    result = subprocess.run([
                        soffice,
                        '--headless',
                        '--invisible',
                        '--nologo',
                        '--nofirststartwizard',
                        '--norestore',
                        '--convert-to', 'pdf:impress_pdf_Export',
                        '--outdir', tmpdir,
                        pptx_path
                    ], capture_output=True, timeout=600, text=True, env=env)
                    
                    print(f"[LibreOffice] return code: {result.returncode}")
                    if result.stderr.strip():
                        print(f"[LibreOffice] stderr: {result.stderr[:300]}")
                    
                    pdf_path = os.path.join(tmpdir, 'input.pdf')
                    if os.path.exists(pdf_path):
                        with open(pdf_path, 'rb') as f:
                            pdf_data = f.read()
                        
                        print(f"[LibreOffice] PDF created, size: {len(pdf_data)} bytes")
                        
                        if len(pdf_data) < 500:
                            raise ConversionError("PDF output is too small")
                        
                        # Thorough PDF validation
                        pdf_doc = fitz.open(stream=pdf_data, filetype="pdf")
                        actual_pages = len(pdf_doc)
                        print(f"[LibreOffice] PDF has {actual_pages} pages")
                        
                        # Check each page for content
                        blank_pages = []
                        low_content_pages = []
                        
                        for page_num in range(actual_pages):
                            page = pdf_doc[page_num]
                            
                            # Get all content indicators
                            text = page.get_text().strip()
                            images = page.get_images()
                            drawings = page.get_drawings()
                            
                            # Also check for vector graphics via display list
                            dl = page.get_displaylist()
                            has_graphics = dl is not None
                            
                            # Get pixmap to check if page is visually blank
                            pix = page.get_pixmap(matrix=fitz.Matrix(0.5, 0.5))  # Low res for speed
                            
                            # Check if page is nearly all white
                            samples = pix.samples
                            total_pixels = pix.width * pix.height
                            white_threshold = 250
                            
                            # Count non-white pixels
                            non_white = 0
                            for i in range(0, len(samples), 3):  # RGB
                                r, g, b = samples[i], samples[i+1], samples[i+2]
                                if r < white_threshold or g < white_threshold or b < white_threshold:
                                    non_white += 1
                            
                            content_ratio = non_white / total_pixels if total_pixels > 0 else 0
                            
                            # Page is blank if it has no text, no images, no drawings, 
                            # and is >99% white pixels
                            is_blank = (
                                not text and 
                                not images and 
                                not drawings and 
                                content_ratio < 0.01
                            )
                            
                            # Page has low content if it should have content but doesn't
                            has_low_content = (
                                slide_contents and 
                                page_num < len(slide_contents) and
                                slide_contents[page_num]['has_content'] and
                                content_ratio < 0.02
                            )
                            
                            if is_blank:
                                blank_pages.append(page_num + 1)
                            elif has_low_content:
                                low_content_pages.append(page_num + 1)
                        
                        pdf_doc.close()
                        
                        # Report findings
                        if blank_pages:
                            print(f"[LibreOffice] BLANK pages: {blank_pages}")
                        if low_content_pages:
                            print(f"[LibreOffice] LOW CONTENT pages (may be missing elements): {low_content_pages}")
                        
                        if expected_pages > 0 and actual_pages != expected_pages:
                            print(f"[LibreOffice] PAGE COUNT MISMATCH: expected {expected_pages}, got {actual_pages}")
                        
                        # In strict mode, fail if ANY issues detected
                        if strict:
                            issues = []
                            if blank_pages:
                                issues.append(f"blank pages: {blank_pages}")
                            if low_content_pages:
                                issues.append(f"low content pages: {low_content_pages}")
                            if expected_pages > 0 and actual_pages != expected_pages:
                                issues.append(f"page count mismatch ({actual_pages} vs {expected_pages})")
                            
                            if issues:
                                raise ConversionError(f"Quality check failed: {'; '.join(issues)}")
                        
                        # Combine blank and low content pages for fixing
                        all_problem_pages = sorted(set(blank_pages + low_content_pages))
                        
                        print(f"[LibreOffice] Success! {actual_pages} pages, {len(all_problem_pages)} need fixing.")
                        
                        if return_blank_info:
                            return pdf_data, all_problem_pages
                        return pdf_data
                    else:
                        print(f"[LibreOffice] PDF not created")
                        files = os.listdir(tmpdir)
                        print(f"[LibreOffice] Files in tmpdir: {files}")
                        
                except FileNotFoundError:
                    continue
                except subprocess.TimeoutExpired:
                    raise ConversionError("LibreOffice conversion timed out (10 min limit)")
                except ConversionError:
                    raise
                except Exception as e:
                    print(f"[LibreOffice] Error with {soffice}: {e}")
                    continue
            
            raise ConversionError("LibreOffice not found or conversion failed")
    
    @staticmethod
    def _convert_basic(buffer: bytes) -> bytes:
        """
        Basic fallback: render slides including images, shapes, and text to PDF.
        This method extracts embedded images and renders them properly.
        """
        if not PPTX_AVAILABLE:
            raise ConversionError("python-pptx library not available")
        
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            from pptx.dml.color import RGBColor
            from pptx.enum.dml import MSO_THEME_COLOR
            from reportlab.lib.utils import ImageReader
            
            prs = Presentation(io.BytesIO(buffer))
            
            # Get slide dimensions
            slide_width = prs.slide_width.pt
            slide_height = prs.slide_height.pt
            
            output = io.BytesIO()
            c = canvas.Canvas(output, pagesize=(slide_width, slide_height))
            
            for slide_num, slide in enumerate(prs.slides, 1):
                # Draw white background first
                c.setFillColorRGB(1, 1, 1)
                c.rect(0, 0, slide_width, slide_height, fill=1)
                
                # Try to render slide background
                try:
                    background = slide.background
                    if background and background.fill:
                        fill = background.fill
                        if fill.type is not None:
                            # Solid color background
                            if hasattr(fill, 'fore_color') and fill.fore_color:
                                try:
                                    rgb = fill.fore_color.rgb
                                    if rgb:
                                        c.setFillColorRGB(rgb[0]/255, rgb[1]/255, rgb[2]/255)
                                        c.rect(0, 0, slide_width, slide_height, fill=1)
                                except:
                                    pass
                except:
                    pass
                
                # Collect all shapes for proper z-order rendering
                shapes_to_render = []
                
                for shape in slide.shapes:
                    try:
                        # Get shape position and size
                        x = shape.left.pt if hasattr(shape, 'left') and hasattr(shape.left, 'pt') else 0
                        y = shape.top.pt if hasattr(shape, 'top') and hasattr(shape.top, 'pt') else 0
                        width = shape.width.pt if hasattr(shape, 'width') and hasattr(shape.width, 'pt') else 100
                        height = shape.height.pt if hasattr(shape, 'height') and hasattr(shape.height, 'pt') else 100
                        
                        # Convert to PDF coordinates (origin at bottom-left)
                        pdf_y = slide_height - y - height
                        
                        shape_info = {
                            'shape': shape,
                            'x': x,
                            'y': pdf_y,
                            'width': width,
                            'height': height,
                            'original_y': y,  # for z-ordering
                        }
                        shapes_to_render.append(shape_info)
                    except Exception as e:
                        continue
                
                # Render shapes (images first, then other shapes, then text on top)
                # First pass: render images and filled shapes
                for shape_info in shapes_to_render:
                    shape = shape_info['shape']
                    x = shape_info['x']
                    y = shape_info['y']
                    width = shape_info['width']
                    height = shape_info['height']
                    
                    try:
                        # Handle pictures/images
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            try:
                                image_stream = io.BytesIO(shape.image.blob)
                                img = Image.open(image_stream)
                                
                                # Convert to RGB if necessary (handles RGBA, P mode, etc.)
                                if img.mode in ('RGBA', 'LA', 'P'):
                                    # Create white background for transparency
                                    background = Image.new('RGB', img.size, (255, 255, 255))
                                    if img.mode == 'P':
                                        img = img.convert('RGBA')
                                    if img.mode in ('RGBA', 'LA'):
                                        background.paste(img, mask=img.split()[-1])
                                        img = background
                                    else:
                                        img = img.convert('RGB')
                                elif img.mode != 'RGB':
                                    img = img.convert('RGB')
                                
                                # Save to buffer for ReportLab
                                img_buffer = io.BytesIO()
                                img.save(img_buffer, format='PNG')
                                img_buffer.seek(0)
                                
                                c.drawImage(ImageReader(img_buffer), x, y, width=width, height=height, preserveAspectRatio=True, mask='auto')
                            except Exception as img_error:
                                # Draw placeholder rectangle for failed images
                                c.setStrokeColorRGB(0.8, 0.8, 0.8)
                                c.setFillColorRGB(0.95, 0.95, 0.95)
                                c.rect(x, y, width, height, fill=1, stroke=1)
                            continue
                        
                        # Handle shapes with fills (rectangles, ovals, etc.)
                        if hasattr(shape, 'fill') and shape.fill:
                            try:
                                fill = shape.fill
                                if fill.type is not None:
                                    # Try to get fill color
                                    if hasattr(fill, 'fore_color') and fill.fore_color:
                                        try:
                                            rgb = fill.fore_color.rgb
                                            if rgb:
                                                c.setFillColorRGB(rgb[0]/255, rgb[1]/255, rgb[2]/255)
                                                c.rect(x, y, width, height, fill=1, stroke=0)
                                        except:
                                            pass
                            except:
                                pass
                        
                        # Handle shapes with lines/borders
                        if hasattr(shape, 'line') and shape.line:
                            try:
                                line = shape.line
                                if line.color and line.color.rgb:
                                    rgb = line.color.rgb
                                    c.setStrokeColorRGB(rgb[0]/255, rgb[1]/255, rgb[2]/255)
                                    line_width = line.width.pt if hasattr(line, 'width') and line.width else 1
                                    c.setLineWidth(line_width)
                                    c.rect(x, y, width, height, fill=0, stroke=1)
                            except:
                                pass
                                
                    except Exception as shape_error:
                        continue
                
                # Second pass: render text on top
                shapes_data = []
                
                for shape_info in shapes_to_render:
                    shape = shape_info['shape']
                    x = shape_info['x']
                    width = shape_info['width']
                    original_y = shape_info['original_y']
                    
                    try:
                        if hasattr(shape, 'left') and hasattr(shape, 'top'):
                            x = shape.left.pt if hasattr(shape.left, 'pt') else 50
                            y = slide_height - (shape.top.pt if hasattr(shape.top, 'pt') else 50)
                            width = shape.width.pt if hasattr(shape, 'width') and hasattr(shape.width, 'pt') else 400
                            
                            # Skip pictures (already rendered)
                            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                                continue
                            
                            if hasattr(shape, "text") and shape.text.strip():
                                # Try to get text color
                                text_color = (0, 0, 0)  # Default black
                                font_bold = False
                                custom_font_size = None
                                
                                try:
                                    if hasattr(shape, 'text_frame'):
                                        for paragraph in shape.text_frame.paragraphs:
                                            for run in paragraph.runs:
                                                if run.font.color and run.font.color.rgb:
                                                    rgb = run.font.color.rgb
                                                    text_color = (rgb[0]/255, rgb[1]/255, rgb[2]/255)
                                                if run.font.bold:
                                                    font_bold = True
                                                if run.font.size:
                                                    custom_font_size = run.font.size.pt
                                                break
                                            break
                                except:
                                    pass
                                
                                shapes_data.append({
                                    'x': x,
                                    'y': y,
                                    'text': shape.text.strip(),
                                    'width': width,
                                    'color': text_color,
                                    'bold': font_bold,
                                    'font_size': custom_font_size
                                })
                    except:
                        pass
                
                # Sort by vertical position (top to bottom)
                shapes_data.sort(key=lambda s: -s['y'])
                
                for shape_data in shapes_data:
                    text = shape_data['text']
                    x = max(30, shape_data['x'])
                    y = shape_data['y']
                    max_width = shape_data['width']
                    text_color = shape_data.get('color', (0, 0, 0))
                    is_bold = shape_data.get('bold', False)
                    custom_font_size = shape_data.get('font_size')
                    
                    # Set text color
                    c.setFillColorRGB(*text_color)
                    
                    # Determine font size based on position (titles are usually at top) or use custom
                    if custom_font_size:
                        font_size = min(custom_font_size, 48)  # Cap at reasonable size
                    elif y > slide_height - 100:
                        font_size = 24
                    else:
                        font_size = 14
                    
                    # Set font
                    font_name = "Helvetica-Bold" if is_bold or y > slide_height - 100 else "Helvetica"
                    c.setFont(font_name, font_size)
                    
                    # Word wrap
                    words = text.split()
                    lines = []
                    current_line = []
                    
                    for word in words:
                        current_line.append(word)
                        test_line = ' '.join(current_line)
                        if c.stringWidth(test_line, font_name, font_size) > max_width - 20:
                            current_line.pop()
                            if current_line:
                                lines.append(' '.join(current_line))
                            current_line = [word]
                    
                    if current_line:
                        lines.append(' '.join(current_line))
                    
                    for line in lines:
                        if y > 30:
                            c.drawString(x, y, line)
                            y -= font_size + 4
                
                # Add slide number at bottom
                c.setFont("Helvetica", 10)
                c.setFillColorRGB(0.5, 0.5, 0.5)
                c.drawString(slide_width - 50, 20, str(slide_num))
                
                c.showPage()
            
            c.save()
            return output.getvalue()
            
        except Exception as e:
            raise ConversionError(f"PowerPoint to PDF conversion failed: {e}")


class HTMLToPDFConverter:
    """Convert HTML/URL to PDF."""
    
    @staticmethod
    async def convert_url(url: str) -> bytes:
        """Convert a webpage URL to PDF using pyppeteer."""
        try:
            from pyppeteer import launch
            import os
            import asyncio
            
            # Use system Chromium if available (Docker)
            chromium_path = os.environ.get('PUPPETEER_EXECUTABLE_PATH', None)
            
            browser = await launch(
                headless=True,
                executablePath=chromium_path,
                args=['--no-sandbox', '--disable-setuid-sandbox', '--disable-dev-shm-usage']
            )
            page = await browser.newPage()
            
            # Set viewport to capture full page
            await page.setViewport({'width': 1920, 'height': 1080})
            
            # Navigate to page
            await page.goto(url, waitUntil='networkidle0', timeout=60000)
            
            # Scroll through entire page to trigger lazy loading
            await page.evaluate('''async () => {
                await new Promise((resolve) => {
                    let totalHeight = 0;
                    const distance = 500;
                    const timer = setInterval(() => {
                        const scrollHeight = document.body.scrollHeight;
                        window.scrollBy(0, distance);
                        totalHeight += distance;
                        if (totalHeight >= scrollHeight) {
                            clearInterval(timer);
                            window.scrollTo(0, 0);
                            resolve();
                        }
                    }, 100);
                });
            }''')
            
            # Wait for images to load
            await page.evaluate('''async () => {
                const images = document.querySelectorAll('img');
                await Promise.all(
                    Array.from(images).map(img => {
                        if (img.complete) return Promise.resolve();
                        return new Promise((resolve) => {
                            img.addEventListener('load', resolve);
                            img.addEventListener('error', resolve);
                            setTimeout(resolve, 5000);
                        });
                    })
                );
            }''')
            
            # Extra wait for any animations/transitions
            await asyncio.sleep(2)
            
            pdf_bytes = await page.pdf({
                'format': 'A4',
                'printBackground': True,
                'margin': {
                    'top': '1cm',
                    'bottom': '1cm',
                    'left': '1cm',
                    'right': '1cm'
                }
            })
            
            await browser.close()
            return pdf_bytes
            
        except ImportError:
            raise ConversionError("pyppeteer library not available")
        except Exception as e:
            raise ConversionError(f"HTML to PDF conversion failed: {e}")
    
    @staticmethod
    def convert_html_string(html: str) -> bytes:
        """Convert HTML string to PDF using reportlab."""
        try:
            from reportlab.platypus import SimpleDocTemplate, Paragraph
            from reportlab.lib.styles import getSampleStyleSheet
            
            output = io.BytesIO()
            doc = SimpleDocTemplate(output, pagesize=letter)
            styles = getSampleStyleSheet()
            
            # Simple HTML to text conversion
            import re
            text = re.sub(r'<[^>]+>', ' ', html)
            text = ' '.join(text.split())
            
            story = [Paragraph(text, styles['Normal'])]
            doc.build(story)
            
            return output.getvalue()
            
        except Exception as e:
            raise ConversionError(f"HTML to PDF conversion failed: {e}")
